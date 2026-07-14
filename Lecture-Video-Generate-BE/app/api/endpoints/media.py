from fastapi import APIRouter, UploadFile, File, HTTPException, Request
from fastapi.responses import JSONResponse
from typing import List
from pathlib import Path
import uuid
import shutil
import subprocess
import requests
import os
import logging
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io
from app.services.gen_slides import SlideGeneratorService


logging.basicConfig(level=logging.INFO)

# Chỉ giữ lại FFMPEG (cần thiết cho video processing)
FFMPEG_CMD = shutil.which("ffmpeg") or shutil.which("ffmpeg.exe") or None

def ensure_executable(cmd: str, name: str):
    if not cmd:
        raise HTTPException(status_code=500, detail=f"'{name}' không tìm thấy trên hệ thống. Vui lòng cài đặt và thêm vào PATH.")

router = APIRouter()

# Thư mục static
ROOT_DIR = Path(__file__).resolve().parents[3]
STATIC_DIR = ROOT_DIR / "static"
PPT_IMAGES_DIR = STATIC_DIR / "ppt_images"
OUTPUTS_DIR = STATIC_DIR / "outputs"
UPLOADS_DIR = ROOT_DIR / "uploads"

for d in (PPT_IMAGES_DIR, OUTPUTS_DIR, UPLOADS_DIR):
    d.mkdir(parents=True, exist_ok=True)

def save_upload_file(tmp_path: Path, upload_file: UploadFile):
    with tmp_path.open("wb") as buffer:
        shutil.copyfileobj(upload_file.file, buffer)

def download_to_file(url: str, dst: Path):
    resp = requests.get(url, stream=True, timeout=60)
    resp.raise_for_status()
    with dst.open("wb") as f:
        for chunk in resp.iter_content(chunk_size=8192):
            if chunk:
                f.write(chunk)

def extract_text_from_shape(shape):
    """
    Extract text từ shape, bao gồm cả table và group shapes
    """
    texts = []

    # Nếu là shape chứa text
    if hasattr(shape, "text_frame") and shape.text_frame:
        for paragraph in shape.text_frame.paragraphs:
            part = "".join(run.text for run in paragraph.runs)
            if part.strip():
                texts.append(part.strip())

    # Nếu là bảng (MSO_SHAPE_TYPE.TABLE = 19)
    if hasattr(shape, 'shape_type') and shape.shape_type == 19:
        try:
            table = shape.table
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        texts.append(cell.text.strip())
        except Exception as e:
            logging.warning(f"Error extracting table text: {e}")

    # Nếu là group shape → duyệt đệ quy
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            texts.extend(extract_text_from_shape(sub_shape))

    return texts

def extract_text_from_slide_ordered(slide):
    """
    Extract text từ slide theo thứ tự từ trên xuống dưới
    """
    all_shapes = []
    
    def collect_shapes(shapes, parent_top=0, parent_left=0):
        """Recursively collect all shapes with absolute positions"""
        for shape in shapes:
            # Tính vị trí tuyệt đối
            abs_top = parent_top + (shape.top if hasattr(shape, 'top') else 0)
            abs_left = parent_left + (shape.left if hasattr(shape, 'left') else 0)
            
            # Nếu là group shape, duyệt đệ quy
            if hasattr(shape, "shapes"):
                collect_shapes(shape.shapes, abs_top, abs_left)
            else:
                # Lưu shape với vị trí tuyệt đối
                all_shapes.append({
                    'shape': shape,
                    'top': abs_top,
                    'left': abs_left
                })
    
    # Collect tất cả shapes
    collect_shapes(slide.shapes)
    
    # Sắp xếp theo top (y-coordinate), sau đó theo left (x-coordinate)
    all_shapes.sort(key=lambda s: (s['top'], s['left']))
    
    # Extract text theo thứ tự
    texts = []
    for item in all_shapes:
        shape = item['shape']
        
        # Extract text từ text_frame
        if hasattr(shape, "text_frame") and shape.text_frame:
            for paragraph in shape.text_frame.paragraphs:
                part = "".join(run.text for run in paragraph.runs)
                if part.strip():
                    texts.append(part.strip())
        
        # Extract text từ table
        elif hasattr(shape, 'shape_type') and shape.shape_type == 19:
            try:
                table = shape.table
                for row in table.rows:
                    row_texts = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_texts.append(cell.text.strip())
                    if row_texts:
                        texts.append(" | ".join(row_texts))  # Join cells với |
            except Exception as e:
                logging.warning(f"Error extracting table text: {e}")
    
    return texts


@router.post("/upload-pptx")
async def upload_pptx(request: Request, file: UploadFile = File(...)):
    """
    Upload PPTX -> chuyển thành ảnh (chỉ sử dụng PIL)
    """
    logging.info(f"Received PPTX upload: {file.filename}")
    
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Chỉ chấp nhận file .pptx")

    job_id = uuid.uuid4().hex
    out_dir = PPT_IMAGES_DIR / job_id
    out_dir.mkdir(parents=True, exist_ok=True)
    tmp_pptx = UPLOADS_DIR / f"{job_id}.pptx"

    try:
        save_upload_file(tmp_pptx, file)
        logging.info(f"Saved PPTX: {tmp_pptx}")

        # Chỉ sử dụng PIL Enhanced
        logging.info("Using PIL Enhanced conversion...")
        return await convert_with_pil_enhanced(tmp_pptx, out_dir, job_id, request)

    except HTTPException:
        raise
    except Exception as e:
        logging.exception("upload_pptx error")
        raise HTTPException(status_code=500, detail=f"Error: {str(e)}")
    
@router.post("/upload-pptx2")
async def upload_pptx_images_only(request: Request, file: UploadFile = File(...)):
    """Upload PPTX -> chuyển thành ảnh GỐC (không xử lý nội dung)"""
    
    if not file.filename.lower().endswith(".pptx"):
        raise HTTPException(status_code=400, detail="Chỉ chấp nhận file .pptx")

    job_id = uuid.uuid4().hex
    out_dir = PPT_IMAGES_DIR / job_id
    out_dir.mkdir(parents=True, exist_ok=True)
    tmp_pptx = UPLOADS_DIR / f"{job_id}.pptx"

    try:
        save_upload_file(tmp_pptx, file)
        return await convert_pptx_to_images_only(tmp_pptx, out_dir, job_id, request)
    except Exception as e:
        logging.exception("upload_pptx2 error")
        raise HTTPException(status_code=500, detail=f"Error: {str(e)}")

@router.post("/extract-pptx-text")
async def extract_pptx_text(file: UploadFile = File(...)):
    """
    Extract text content từ PPTX theo thứ tự từ trên xuống dưới
    Hỗ trợ: text frames, tables, group shapes
    """
    try:
        # Save uploaded file
        tmp_path = UPLOADS_DIR / f"temp_{file.filename}"
        with open(tmp_path, 'wb') as f:
            content = await file.read()
            f.write(content)
        
        logging.info(f"Extracting text from: {tmp_path}")
        
        # Initialize Service for Rewriting
        slide_service = SlideGeneratorService() 
        
        # Extract text using python-pptx
        prs = Presentation(str(tmp_path))
        
        slides_to_process = []
        raw_slides_data = []

        for slide_idx, slide in enumerate(prs.slides):
            slide_text_parts = extract_text_from_slide_ordered(slide)
            
            title = slide_text_parts[0] if slide_text_parts else f"Slide {slide_idx + 1}"
            
            full_content = "\n".join(slide_text_parts)
            
            raw_data = {
                "slide_number": slide_idx,
                "title": title,
                "content": full_content,
                "all_text": slide_text_parts
            }
            raw_slides_data.append(raw_data)
            
            # Add to batch for rewriting if content is not empty
            if full_content.strip():
                slides_to_process.append({
                    "id": slide_idx,
                    "content": full_content
                })
                
        logging.info(f"Prepared {len(slides_to_process)} slides for batch rewriting")

        # Rewrite content using Gemini Batch
        rewritten_map = slide_service.rewrite_batch_slide_content(slides_to_process)
        
        # Merge results
        slides_text = []
        for raw in raw_slides_data:
            idx = raw["slide_number"]
            # Get rewritten content or fallback to original
            rewritten = rewritten_map.get(idx, raw["content"])
            
            slides_text.append({
                "slide_number": idx,
                "title": raw["title"],
                "content": raw["content"],
                "rewritten_content": rewritten,
                "all_text": raw["all_text"]
            })
            
            logging.info(f"Slide {idx}: Processed")
        
        # Clean up
        tmp_path.unlink(missing_ok=True)
        
        return JSONResponse({
            "success": True,
            "slides_text": slides_text,
            "total_slides": len(slides_text)
        })
        
    except Exception as e:
        logging.error(f"Error extracting PPTX text: {e}")
        raise HTTPException(status_code=500, detail=str(e))

async def convert_with_pil_enhanced(tmp_pptx: Path, out_dir: Path, job_id: str, request: Request):
    """Convert PPTX using python-pptx + PIL - ENHANCED AND FIXED"""
    
    try:
        prs = Presentation(str(tmp_pptx))
        total_slides = len(prs.slides)
        logging.info(f"PIL Enhanced: Processing {total_slides} slides")
        
        slides_data = []
        
        # Get font function
        def get_font(size):
            font_paths = [
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
                "C:/Windows/Fonts/arial.ttf"
            ]
            
            for font_path in font_paths:
                if Path(font_path).exists():
                    try:
                        return ImageFont.truetype(font_path, size)
                    except:
                        continue
            
            return ImageFont.load_default()
        
        for slide_idx, slide in enumerate(prs.slides):
            logging.info(f"Processing slide {slide_idx + 1}/{total_slides}")
            
            # Create high-quality image
            img = Image.new('RGB', (1920, 1080), color='white')
            draw = ImageDraw.Draw(img)
            
            # Extract content
            title_text = ""
            content_texts = []
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    
                    # Classify by position
                    if hasattr(shape, 'top') and shape.top < 200000:  # Upper area
                        if not title_text:
                            title_text = text
                        else:
                            content_texts.append(text)
                    else:
                        content_texts.append(text)
            
            # Fallback title
            if not title_text and content_texts:
                title_text = content_texts.pop(0)
            
            if not title_text:
                title_text = f"Slide {slide_idx + 1}"
            
            # Setup fonts
            title_font = get_font(48)
            content_font = get_font(32)
            small_font = get_font(24)
            
            # Draw header
            draw.rectangle([(0, 0), (1920, 100)], fill='#2c3e50')
            draw.text((1850, 50), f"Slide {slide_idx + 1}", fill='white', 
                     font=small_font, anchor="rm")
            
            y_pos = 140
            
            # Draw title
            title_lines = smart_wrap_text(draw, title_text, title_font, 1700)
            for line in title_lines:
                draw.text((100, y_pos), line, fill='#2c3e50', font=title_font)
                y_pos += 60
            
            # Draw underline
            draw.line([(100, y_pos + 10), (1820, y_pos + 10)], fill='#3498db', width=3)
            y_pos += 50
            
            # Draw content
            for content in content_texts:
                if '\n' in content:
                    # Multi-line content
                    lines = content.split('\n')
                    for line in lines:
                        if line.strip():
                            formatted_line = format_bullet_point(line.strip())
                            content_lines = smart_wrap_text(draw, formatted_line, content_font, 1600)
                            for content_line in content_lines:
                                draw.text((150, y_pos), content_line, fill='#2c3e50', font=content_font)
                                y_pos += 40
                            y_pos += 10
                else:
                    # Single line
                    formatted_content = format_bullet_point(content)
                    content_lines = smart_wrap_text(draw, formatted_content, content_font, 1600)
                    for content_line in content_lines:
                        draw.text((150, y_pos), content_line, fill='#2c3e50', font=content_font)
                        y_pos += 40
                    y_pos += 15
            
            # Save image
            slide_filename = f"slide_{slide_idx}.png"
            slide_path = out_dir / slide_filename
            img.save(slide_path, 'PNG', quality=95)
            
            base = str(request.base_url).rstrip("/")
            image_url = f"{base}/static/ppt_images/{job_id}/{slide_filename}"
            
            slides_data.append({
                "id": f"{job_id}_{slide_idx}",
                "image_url": image_url,
                "order": slide_idx,
                "slide_number": slide_idx,
                "title": title_text,
                "content": content_texts
            })
            
            logging.info(f"PIL created slide {slide_idx}: {slide_filename}")
        
        pptx_url = f"{str(request.base_url).rstrip('/')}/static/uploads/{tmp_pptx.name}"
        
        logging.info(f"PIL successfully converted {len(slides_data)} slides")
        return JSONResponse({"success": True, "pptx_url": pptx_url, "slides": slides_data})
        
    except Exception as e:
        logging.error(f"PIL Enhanced conversion failed: {e}")
        raise Exception(f"PIL Enhanced error: {str(e)}")

def smart_wrap_text(draw, text, font, max_width):
    """Smart text wrapping"""
    if not text:
        return []
        
    words = text.split()
    lines = []
    current_line = []
    
    for word in words:
        test_line = ' '.join(current_line + [word])
        try:
            bbox = draw.textbbox((0, 0), test_line, font=font)
            width = bbox[2] - bbox[0]
        except:
            # Fallback for older PIL
            width = draw.textsize(test_line, font=font)[0]
        
        if width <= max_width:
            current_line.append(word)
        else:
            if current_line:
                lines.append(' '.join(current_line))
                current_line = [word]
            else:
                lines.append(word)
    
    if current_line:
        lines.append(' '.join(current_line))
    
    return lines

def format_bullet_point(text):
    """Add bullet point"""
    text = text.strip()
    if not text.startswith('•') and not text.startswith('-') and not text.startswith('*'):
        return f"• {text}"
    return text

@router.post("/combine-slide")
async def combine_slide(request: Request, data: dict):
    """
    Ghép ảnh slide với video -> video đặt góc dưới bên PHẢI
    """
    image_url = data.get("image_url")
    video_url = data.get("video_url")
    if not image_url or not video_url:
        raise HTTPException(status_code=400, detail="Thiếu image_url hoặc video_url")

    job_id = uuid.uuid4().hex
    out_file = OUTPUTS_DIR / f"slide_compose_{job_id}.mp4"

    try:
        # Lấy file image và video (giống như cũ)
        if str(request.base_url).rstrip("/") in image_url and "/static/" in image_url:
            rel = image_url.split("/static/")[-1]
            img_path = STATIC_DIR / Path(rel)
        else:
            img_path = UPLOADS_DIR / f"{job_id}_img"
            download_to_file(image_url, img_path)

        if str(request.base_url).rstrip("/") in video_url and "/static/" in video_url:
            rel = video_url.split("/static/")[-1]
            vid_path = STATIC_DIR / Path(rel)
        else:
            vid_path = UPLOADS_DIR / f"{job_id}_vid"
            download_to_file(video_url, vid_path)

        # Tạo background video từ ảnh
        tmp_bg_video = UPLOADS_DIR / f"{job_id}_bg.mp4"
        ensure_executable(FFMPEG_CMD, "ffmpeg")
        
        cmd_bg = [
            FFMPEG_CMD, "-y",
            "-loop", "1",
            "-i", str(img_path),
            "-c:v", "libx264",
            "-t", "1",
            "-vf", "scale=1920:1080",
            "-pix_fmt", "yuv420p",
            str(tmp_bg_video)
        ]
        proc_bg = subprocess.run(cmd_bg, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=60)
        
        # Overlay video ở góc dưới bên PHẢI với padding 20px
        overlay_scale = "400:-1"  # Width 400px, height auto
        cmd_overlay = [
            FFMPEG_CMD, "-y",
            "-i", str(tmp_bg_video),
            "-i", str(vid_path),
            "-filter_complex", 
            f"[1:v]scale={overlay_scale}[ov];[0:v][ov]overlay=main_w-overlay_w-20:main_h-overlay_h-20",
            "-c:v", "libx264",
            "-c:a", "aac",
            "-shortest",
            str(out_file)
        ]
        proc_overlay = subprocess.run(cmd_overlay, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=300)

        result_url = f"{str(request.base_url).rstrip('/')}/static/outputs/{out_file.name}"
        return JSONResponse({"result_url": result_url})
        
    except subprocess.CalledProcessError as e:
        stderr = e.stderr.decode(errors="ignore") if e.stderr else str(e)
        logging.error("FFmpeg error: %s", stderr)
        raise HTTPException(status_code=500, detail=f"FFmpeg error: {stderr}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/concat-videos")
async def concat_videos(request: Request, data: dict):
    """
    Nối nhiều video liền nhau (không có delay)
    """
    videos = data.get("videos")
    if not videos or not isinstance(videos, list):
        raise HTTPException(status_code=400, detail="videos phải là mảng URL")

    job_id = uuid.uuid4().hex
    tmp_dir = UPLOADS_DIR / job_id
    tmp_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        # Download tất cả videos
        local_paths = []
        for idx, vurl in enumerate(videos):
            if str(request.base_url).rstrip("/") in vurl and "/static/" in vurl:
                rel = vurl.split("/static/")[-1]
                local = STATIC_DIR / Path(rel)
            else:
                local = tmp_dir / f"part_{idx}.mp4"
                download_to_file(vurl, local)
            if not local.exists():
                raise HTTPException(status_code=500, detail=f"Không tìm thấy video: {vurl}")
            local_paths.append(local)

        out_file = OUTPUTS_DIR / f"concat_{job_id}.mp4"
        ensure_executable(FFMPEG_CMD, "ffmpeg")
        
        if len(local_paths) == 1:
            # Chỉ 1 video, copy trực tiếp
            cmd = [
                FFMPEG_CMD, "-y",
                "-i", str(local_paths[0]),
                "-c", "copy",
                str(out_file)
            ]
            proc_concat = subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=600)
        else:
            # Nhiều videos: Tạo file list và concat trực tiếp
            concat_list = tmp_dir / "concat_list.txt"
            with open(concat_list, 'w', encoding='utf-8') as f:
                for video_path in local_paths:
                    # Sử dụng đường dẫn tuyệt đối và escape ký tự đặc biệt
                    safe_path = str(video_path).replace('\\', '/')
                    f.write(f"file '{safe_path}'\n")
            
            # Log nội dung file với encoding an toàn
            try:
                with open(concat_list, 'r', encoding='utf-8') as f:
                    concat_content = f.read()
                logging.info(f"Concat list content:\n{concat_content}")
            except UnicodeDecodeError:
                logging.warning("Cannot read concat list content due to encoding issues")
            
            # Concat sử dụng concat demuxer (không có delay)
            cmd = [
                FFMPEG_CMD, "-y",
                "-f", "concat",
                "-safe", "0",
                "-i", str(concat_list),
                "-c", "copy",
                str(out_file)
            ]
            
            logging.info(f"Running concat command: {' '.join(cmd)}")
            proc_concat = subprocess.run(cmd, check=True, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=600)
        
        result_url = f"{str(request.base_url).rstrip('/')}/static/outputs/{out_file.name}"
        return JSONResponse({"result_url": result_url})
        
    except subprocess.CalledProcessError as e:
        stderr = e.stderr.decode(errors="ignore") if e.stderr else str(e)
        logging.error("FFmpeg concat error: %s", stderr)
        raise HTTPException(status_code=500, detail=f"FFmpeg concat error: {stderr}")
    except Exception as e:
        logging.exception("concat_videos unexpected error")
        raise HTTPException(status_code=500, detail=str(e))

async def convert_pptx_to_images_only(tmp_pptx: Path, out_dir: Path, job_id: str, request: Request):
    """
    Convert PPTX to images - GIỮ NGUYÊN SLIDE GỐC
    Sử dụng LibreOffice để export ảnh chất lượng cao
    """
    import subprocess
    import shutil
    import os
    
    try:
        # Method 1: PowerPoint COM (Windows only)
        if os.name == 'nt':  
            try:
                import comtypes.client
                import pythoncom
                
                logging.info("Using PowerPoint COM to export slides as images...")
                pythoncom.CoInitialize()
                powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
                powerpoint.Visible = 0
                
                presentation = powerpoint.Presentations.Open(str(tmp_pptx.absolute()), WithWindow=False)
                
                slides_data = []
                total_slides = presentation.Slides.Count
                
                for slide_idx in range(1, total_slides + 1):
                    slide = presentation.Slides(slide_idx)
                    
                    slide_filename = f"slide_{slide_idx - 1}.png"
                    slide_path = out_dir / slide_filename
                    
                    slide.Export(str(slide_path.absolute()), "PNG", 1920, 1080)
                    
                    base = str(request.base_url).rstrip("/")
                    image_url = f"{base}/static/ppt_images/{job_id}/{slide_filename}"
                    
                    slides_data.append({
                        "id": f"{job_id}_{slide_idx - 1}",
                        "image_url": image_url,
                        "order": slide_idx - 1,
                        "slide_number": slide_idx - 1,
                    })
                    
                    logging.info(f"Exported slide {slide_idx}/{total_slides}: {slide_filename}")
                
                presentation.Close()
                powerpoint.Quit()
                pythoncom.CoUninitialize()
                
                pptx_url = f"{str(request.base_url).rstrip('/')}/static/uploads/{tmp_pptx.name}"
                
                logging.info(f"PowerPoint COM successfully exported {len(slides_data)} slides")
                return JSONResponse({"success": True, "pptx_url": pptx_url, "slides": slides_data})
                
            except ImportError:
                logging.warning("pywin32/comtypes not available on this system")
            except Exception as e:
                logging.error(f"PowerPoint COM failed: {e}")
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass
        
        # Method 2: LibreOffice (Linux/Mac/Windows)
        logging.info("Using LibreOffice conversion...")
        
        # Check defaults first, then user provided path
        libreoffice_cmd = shutil.which("soffice") or shutil.which("libreoffice")
        
        # User defined path fallback
        if not libreoffice_cmd and os.name == 'nt':
             manual_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
             if os.path.exists(manual_path):
                 libreoffice_cmd = manual_path

        if not libreoffice_cmd:
            raise Exception("LibreOffice not found. Please install: sudo apt install libreoffice")
        
        # Convert PPTX to PDF first
        pdf_dir = UPLOADS_DIR / f"{job_id}_pdf"
        pdf_dir.mkdir(exist_ok=True)
        
        cmd_pdf = [
            libreoffice_cmd,
            "--headless",
            "--convert-to", "pdf",
            "--outdir", str(pdf_dir),
            str(tmp_pptx.absolute())
        ]
        
        logging.info(f"Running LibreOffice: {' '.join(cmd_pdf)}")
        subprocess.run(cmd_pdf, check=True, timeout=120)
        
        # Find the generated PDF
        pdf_files = list(pdf_dir.glob("*.pdf"))
        if not pdf_files:
            raise Exception("LibreOffice did not generate PDF")
        
        pdf_path = pdf_files[0]
        logging.info(f"PDF generated: {pdf_path}")
        
        # Convert PDF to images using pdf2image
        try:
            from pdf2image import convert_from_path
            
            logging.info("Converting PDF to images...")
            
            # User defined poppler path
            poppler_path = None
            if os.name == 'nt':
                manual_poppler = r"D:\Downloads\poppler\poppler-25.12.0\Library\bin"
                if os.path.exists(manual_poppler):
                    poppler_path = manual_poppler

            images = convert_from_path(str(pdf_path), dpi=200, poppler_path=poppler_path)
            
            slides_data = []
            for slide_idx, img in enumerate(images):
                slide_filename = f"slide_{slide_idx}.png"
                slide_path = out_dir / slide_filename
                
                # Resize to standard size and save
                img_resized = img.resize((1920, 1080), Image.Resampling.LANCZOS)
                img_resized.save(slide_path, 'PNG', quality=95, optimize=True)
                
                base = str(request.base_url).rstrip("/")
                image_url = f"{base}/static/ppt_images/{job_id}/{slide_filename}"
                
                slides_data.append({
                    "id": f"{job_id}_{slide_idx}",
                    "image_url": image_url,
                    "order": slide_idx,
                    "slide_number": slide_idx,
                })
                
                logging.info(f"Converted slide {slide_idx + 1}: {slide_filename}")
            
            # Clean up PDF
            shutil.rmtree(pdf_dir, ignore_errors=True)
            
            pptx_url = f"{str(request.base_url).rstrip('/')}/static/uploads/{tmp_pptx.name}"
            
            logging.info(f"LibreOffice + pdf2image successfully converted {len(slides_data)} slides")
            return JSONResponse({"success": True, "pptx_url": pptx_url, "slides": slides_data})
            
        except ImportError:
            logging.error("pdf2image not installed. Install with: pip install pdf2image")
            raise Exception("pdf2image library is required. Run: pip install pdf2image")
        
    except Exception as e:
        logging.error(f"PPTX to image conversion failed: {e}")
        raise HTTPException(status_code=500, detail=f"PPTX conversion error: {str(e)}")