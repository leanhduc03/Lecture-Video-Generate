import google.generativeai as genai
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import json
import os
from typing import Optional, Dict, List
from datetime import datetime
from app.prompts.templates import SLIDE_STRUCTURE_PROMPT, SLIDE_CONTENT_REWRITE_BATCH_PROMPT
from dotenv import load_dotenv

load_dotenv()

class SlideGeneratorService:
    """Service để tạo slides từ nội dung văn bản"""
    
    def __init__(
        self, 
        api_key: str = None,
        rewrite_api_key: str = None,
        model_name: str = None
    ):
        """
        Khởi tạo service
        
        Args:
            api_key: Google Gemini API key cho tạo slide
            rewrite_api_key: Google Gemini API key cho rewrite content (optional)
            model_name: Tên model Gemini AI
        """
        self.gen_api_key = api_key or os.getenv("GEMINI_API_KEY")
        self.rewrite_api_key = rewrite_api_key or os.getenv("GEMINI_REWRITE_API_KEY")
        self.model_name = model_name or os.getenv("GEMINI_MODEL_NAME", "models/gemini-2.5-flash")
        self.model = None

    def _configure_genai(self, api_key: str):
        """Configure GenAI with specific key"""
        genai.configure(api_key=api_key)
        self.model = genai.GenerativeModel(self.model_name)
    
    def generate_slide_structure(self, content: str, num_slides: Optional[int] = None) -> Dict:
        """
        Tạo cấu trúc slide từ nội dung
        
        Args:
            content: Nội dung cần tạo slide
            num_slides: Số lượng slide mong muốn (None = tự động 5-10 slides)
            
        Returns:
            Dict với structure:
            {
                "success": bool,
                "data": {
                    "title": str,
                    "slides": [...]
                },
                "message": str,
                "error": str (nếu có)
            }
        """
        try:
            # Validate input
            if not content or not content.strip():
                return {
                    "success": False,
                    "data": None,
                    "message": "Nội dung trống!",
                    "error": "EMPTY_CONTENT"
                }
            
            # Xác định số lượng slide
            if num_slides and num_slides > 0:
                slide_range = f"{num_slides} slide"
            else:
                slide_range = "5-10 slide"
            
            prompt = SLIDE_STRUCTURE_PROMPT.format(slide_range=slide_range, content=content)
            
            # Configure Key for Generation
            self._configure_genai(self.gen_api_key)
            
            response = self.model.generate_content(prompt)
            
            # Làm sạch response để extract JSON
            json_text = response.text.strip()
            
            # Loại bỏ markdown formatting
            if json_text.startswith('```json'):
                json_text = json_text[7:]
            elif json_text.startswith('```'):
                json_text = json_text[3:]
            
            if json_text.endswith('```'):
                json_text = json_text[:-3]
            
            json_text = json_text.strip()
            
            slide_data = json.loads(json_text)
            
            # Validation: Kiểm tra độ dài original_content
            for slide_info in slide_data.get('slides', []):
                original_content = slide_info.get('original_content', '')
                if len(original_content) > 250:
                    # Tự động cắt nội dung nếu vượt quá 250 ký tự
                    slide_info['original_content'] = original_content[:247] + '...'
            
            return {
                "success": True,
                "data": slide_data,
                "message": f"Tạo thành công {len(slide_data['slides'])} slides",
                "error": None
            }
            
        except json.JSONDecodeError as e:
            return {
                "success": False,
                "data": None,
                "message": "Lỗi parse JSON từ AI",
                "error": f"JSON_DECODE_ERROR: {str(e)}"
            }
        except Exception as e:
            return {
                "success": False,
                "data": None,
                "message": "Lỗi khi gọi API",
                "error": f"API_ERROR: {str(e)}"
            }

    def rewrite_batch_slide_content(self, slides_data: List[Dict]) -> Dict[int, str]:
        """
        Viết lại nội dung slide cho nhiều slide cùng lúc (Batch processing)
        Tự động chia thành các batch nhỏ nếu quá nhiều slides
        
        Args:
            slides_data: List các dict [{"id": int, "content": str}]
            
        Returns:
            Dict mapping {id: rewritten_content}
        """
        try:
            if not slides_data:
                return {}
            
            BATCH_SIZE = 5  # Xử lý 5 slides/lần
            result = {}
            
            for i in range(0, len(slides_data), BATCH_SIZE):
                batch = slides_data[i:i + BATCH_SIZE]
                print(f"Processing batch {i//BATCH_SIZE + 1}: slides {i+1}-{min(i+BATCH_SIZE, len(slides_data))}")
                
                # Configure Key for Rewrite
                self._configure_genai(self.rewrite_api_key)
                
                input_json = json.dumps(batch, ensure_ascii=False)
                prompt = SLIDE_CONTENT_REWRITE_BATCH_PROMPT.format(content=input_json)
                
                # Retry logic
                max_retries = 2
                for attempt in range(max_retries):
                    try:
                        response = self.model.generate_content(
                            prompt
                        )
                        
                        # Extract JSON
                        text = response.text.strip()
                        if text.startswith('```json'):
                            text = text[7:]
                        elif text.startswith('```'):
                            text = text[3:]
                        if text.endswith('```'):
                            text = text[:-3]
                        text = text.strip()
                        
                        rewritten_data = json.loads(text)
                        for item in rewritten_data:
                            result[item['id']] = item['rewritten_content']
                        
                        break  # Success, exit retry loop
                        
                    except Exception as retry_error:
                        print(f"Batch {i//BATCH_SIZE + 1} attempt {attempt + 1}/{max_retries} failed: {repr(retry_error)}")
                        if attempt == max_retries - 1:
                            # Fallback: use original content for this batch
                            for item in batch:
                                result[item['id']] = item['content']
                        else:
                            import time
                            time.sleep(2)
                
                # Delay giữa các batch để tránh rate limit
                if i + BATCH_SIZE < len(slides_data):
                    import time
                    time.sleep(1)
            
            return result
            
        except Exception as e:
            print(f"Error batch rewriting content: {repr(e)}")
            # Fallback: return original content
            return {item['id']: item['content'] for item in slides_data}
    
    def create_presentation(self, slide_data: Dict) -> Presentation:
        """
        Tạo PowerPoint presentation hoàn chỉnh
        
        Args:
            slide_data: Dữ liệu cấu trúc slide
            
        Returns:
            Presentation object
        """
        prs = Presentation()
        
        # Slide tiêu đề
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        
        title.text = slide_data['title']

        # Tạo các slide nội dung
        for slide_info in slide_data['slides']:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            
            # Tiêu đề slide
            title_shape = slide.shapes.title
            title_shape.text = slide_info['title']
            
            # Nội dung slide
            content_shape = slide.placeholders[1]
            text_frame = content_shape.text_frame
            text_frame.clear()
            
            for i, point in enumerate(slide_info['content']):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                
                # Loại bỏ dấu bullet nếu có
                clean_text = point.strip()
                if clean_text.startswith('•'):
                    clean_text = clean_text[1:].strip()
                if clean_text.startswith('-'):
                    clean_text = clean_text[1:].strip()
                
                p.text = clean_text
                p.level = 0
                
                # Định dạng text
                font = p.font
                font.name = 'Arial'
                font.size = Pt(18)
                font.color.rgb = RGBColor(64, 64, 64)
        
        return prs
    
    def save_presentation(self, prs: Presentation, output_dir: str, filename: str = None) -> Dict:
        """
        Lưu PowerPoint presentation
        
        Args:
            prs: Presentation object
            output_dir: Thư mục output
            filename: Tên file (optional, sẽ tự generate nếu không có)
            
        Returns:
            Dict với thông tin file đã lưu
        """
        try:
            # Tạo thư mục nếu chưa có
            os.makedirs(output_dir, exist_ok=True)
            
            # Generate filename nếu không có
            if not filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"presentation_{timestamp}.pptx"
            
            # Đảm bảo có extension .pptx
            if not filename.endswith('.pptx'):
                filename += '.pptx'
            
            filepath = os.path.join(output_dir, filename)
            prs.save(filepath)
            
            # Lấy thông tin file
            file_size = os.path.getsize(filepath)
            
            return {
                "success": True,
                "filepath": filepath,
                "filename": filename,
                "file_size": file_size,
                "message": "Lưu PowerPoint thành công"
            }
            
        except Exception as e:
            return {
                "success": False,
                "filepath": None,
                "filename": None,
                "file_size": 0,
                "message": f"Lỗi khi lưu PowerPoint: {str(e)}"
            }
    
    def save_json(self, slide_data: Dict, output_dir: str, filename: str = None) -> Dict:
        """
        Lưu dữ liệu slide dưới dạng JSON
        
        Args:
            slide_data: Dữ liệu cấu trúc slide
            output_dir: Thư mục output
            filename: Tên file (optional)
            
        Returns:
            Dict với thông tin file đã lưu
        """
        try:
            # Tạo thư mục nếu chưa có
            os.makedirs(output_dir, exist_ok=True)
            
            # Generate filename nếu không có
            if not filename:
                timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
                filename = f"slide_data_{timestamp}.json"
            
            # Đảm bảo có extension .json
            if not filename.endswith('.json'):
                filename += '.json'
            
            filepath = os.path.join(output_dir, filename)
            
            with open(filepath, 'w', encoding='utf-8') as f:
                json.dump(slide_data, f, ensure_ascii=False, indent=2)
            
            file_size = os.path.getsize(filepath)
            
            return {
                "success": True,
                "filepath": filepath,
                "filename": filename,
                "file_size": file_size,
                "message": "Lưu JSON thành công"
            }
            
        except Exception as e:
            return {
                "success": False,
                "filepath": None,
                "filename": None,
                "file_size": 0,
                "message": f"Lỗi khi lưu JSON: {str(e)}"
            }
    
    def generate_slides_complete(
        self, 
        content: str, 
        output_dir: str = "./output/presentations",
        num_slides: Optional[int] = None,
        pptx_filename: Optional[str] = None,
        json_filename: Optional[str] = None
    ) -> Dict:
        """
        Function hoàn chỉnh: Generate slides, lưu PPTX và JSON
        (Không cần convert thành images - sẽ dùng /media/upload-pptx)
        """
        # Step 1: Generate slide structure
        structure_result = self.generate_slide_structure(content, num_slides)
        
        if not structure_result["success"]:
            return structure_result
        
        slide_data = structure_result["data"]
        
        # Step 2: Create PowerPoint
        try:
            prs = self.create_presentation(slide_data)
        except Exception as e:
            return {
                "success": False,
                "data": None,
                "message": "Lỗi khi tạo PowerPoint",
                "error": f"CREATE_PPTX_ERROR: {str(e)}"
            }
        
        # Step 3: Save PowerPoint
        pptx_result = self.save_presentation(prs, output_dir, pptx_filename)
        
        if not pptx_result["success"]:
            return {
                "success": False,
                "data": None,
                "message": pptx_result["message"],
                "error": "SAVE_PPTX_ERROR"
            }
        
        # Step 4: Save JSON
        json_result = self.save_json(slide_data, output_dir, json_filename)
        
        # Return combined result
        return {
            "success": True,
            "data": {
                "slide_data": slide_data,
                "num_slides": len(slide_data['slides']) + 1,  # +1 cho slide tiêu đề
                "pptx_file": {
                    "filepath": pptx_result["filepath"],
                    "filename": pptx_result["filename"],
                    "file_size": pptx_result["file_size"]
                },
                "json_file": {
                    "filepath": json_result["filepath"],
                    "filename": json_result["filename"],
                    "file_size": json_result["file_size"]
                } if json_result["success"] else None
            },
            "message": f"Tạo thành công presentation với {len(slide_data['slides']) + 1} slides"
        }


# UTILITY FUNCTIONS
def get_available_models(api_key: str) -> List[str]:
    """
    Lấy danh sách các model có sẵn
    
    Args:
        api_key: Google API key
        
    Returns:
        List tên model có sẵn
    """
    try:
        genai.configure(api_key=api_key)
        models = genai.list_models()
        available_models = []
        for model in models:
            if 'generateContent' in model.supported_generation_methods:
                available_models.append(model.name)
        return available_models
    except Exception as e:
        print(f"Lỗi khi lấy danh sách model: {e}")
        return []